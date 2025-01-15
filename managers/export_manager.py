import datetime as dt
import pandas as pd
import time
import pandas as pd
from langchain_core.prompts import ChatPromptTemplate
from langchain_anthropic import ChatAnthropic
import datetime as dt
import tqdm
import time
import streamlit as st
import openpyxl
import json, base64, os
from io import BytesIO
import pandas as pd
import urllib3
from urllib3.util.retry import Retry
from requests.adapters import HTTPAdapter
import requests
from time import sleep
import re
import sys
from streamlit_gsheets import GSheetsConnection
import collections 
import collections.abc
import subprocess

import collections 
import collections.abc
import subprocess
import sys
import pandas as pd
import requests
import streamlit as st

try:
    from pptx import Presentation
except ModuleNotFoundError:
    print("Module 'pptx' not found. Installing it now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    print("Module 'pptx' installed successfully.")

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import json, os, io, base64

class ExportManager:

    topic_colors = {
    "social": RGBColor(186, 85, 211),
    "technological": RGBColor(30, 144, 255),
    "environmental": RGBColor(10, 160, 160),
    "economic": RGBColor(220, 20, 60),
    "political": RGBColor(139, 69, 19),
    "business_and_investment": RGBColor(255, 165, 0)
    }

    topic_titles = {
        "social": "social",
        "technological": "technological",
        "environmental": "environmental",
        "economic": "economic",
        "political": "political",
        "business_and_investment": "business_and_investment"
    }

    # *** PPTX Utilities ***
    @staticmethod
    def init_slides():
        prs = Presentation()
        prs.slide_width = Inches(16.0)  
        prs.slide_height = Inches(9.0)  
        return prs

    # Define function to add slides with optional fixed font size
    @staticmethod
    def add_slide(prs, topic, title, content, title_font_size, content_font_size, background_color=RGBColor(240, 240, 240), fixed_font_size=False, rgb_color = "#F3B915"):
        slide_layout = prs.slide_layouts[1]  # Use title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]

        # Adjust title placeholder size to avoid overlapping with the black line
        title_placeholder.width = Inches(12.5)
        title_placeholder.height = Inches(1.2)
        title_placeholder.left = Inches(1)
        title_placeholder.top = Inches(0.4)

        # Set title and font size
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.clear()
        title_p = title_text_frame.add_paragraph()
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.bold = True
        title_p.font.name = 'Microsoft JhengHei'
        title_p.font.color.rgb = ExportManager.topic_colors.get(topic, rgb_color)
        title_p.alignment = PP_ALIGN.LEFT

        # Add a black line below the title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.9), Inches(14), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.width = Inches(0)

        # Add solid circle in the upper-right corner
        solid_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(13.7), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        solid_circle.fill.solid()
        solid_circle.fill.fore_color.rgb = ExportManager.topic_colors.get(topic, rgb_color)

        # Add hollow circle in the upper-right corner
        hollow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(14.5), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        hollow_circle.line.color.rgb = RGBColor(0, 0, 0)
        hollow_circle.line.width = Pt(1)
        hollow_circle.fill.background()

        # Adjust content placeholder size
        body_placeholder.width = Inches(14)
        body_placeholder.height = Inches(6)
        body_placeholder.left = Inches(1)
        body_placeholder.top = Inches(2.05)

        # Set content and font size
        text_frame = body_placeholder.text_frame
        text_frame.clear()

        # Check if font size should be fixed or auto-adjusted
        if fixed_font_size:
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-sizing
        else:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Enable auto-sizing

        paragraphs = content.split('\n')
        for paragraph in paragraphs:
            p = text_frame.add_paragraph()
            parts = paragraph.split('**')
            for i, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(content_font_size)
                run.font.name = 'Microsoft JhengHei'
                if i % 2 == 1:  # Odd index means text between asterisks
                    run.font.bold = True

    # Use this function with `fixed_font_size=True` for Trend Insights and Hashtag slides
    @staticmethod
    def add_report_and_keywords(prs, data, topic, rgb_color):
        # For the Trend Insights and Hashtag slide, use fixed font size
        # todo *************
        slide_title = ExportManager.topic_titles.get(topic, topic)
        ExportManager.add_slide(prs = prs,
                title = slide_title, 
                content = "\n\n".join([f"{key}: {value['<a>趨勢名稱']}" for key, value in data.items()]),
                title_font_size=32, content_font_size=26, fixed_font_size=False, topic = topic, rgb_color = rgb_color)
    
    # Define function to add main trend pages (for STEEP)
    @staticmethod
    def add_trend_pages(prs, data, trend_number, trend_type, topic, rgb_color):
        trend_key = f"{trend_type}{trend_number}"
        if trend_key in data:
            trend = data[trend_key]

            insight = trend['<b>概述']
            events_content = "\n".join(
                [event
                # if isinstance(event, dict) else str(event)
                for event in trend["<d>案例"]]
            )
            subissues = "、".join(str(hashtag) for hashtag in trend["<c>關鍵字"])

            # First Slide: Trend Insights & Important Events & Subissues
            ExportManager.add_slide(prs = prs,
                topic = topic,
                title = f"{trend_type}{trend_number}： {trend['<a>趨勢名稱']}",
                content = f"**趨勢洞察：**{insight}\n\n**代表事件：**\n{events_content}\n\n**子議題：**\n{subissues}"
                , title_font_size=32, content_font_size=22 # Adjust content font size to 26
                , rgb_color = rgb_color
            )

            # Second Slide: Representative Events & Important stakeholders
            key_datas = "\n".join([key_data for key_data in trend['<e>報告中的相關數據點']])
            important_stakeholders = "\n".join([f" {key}: \n   {"\n   ".join(value)}" for key, value in trend["<i>重要關係人"].items()])

            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}： {trend['<a>趨勢名稱']}",
                content = f"**關鍵數據：**\n{key_datas}\n\n**重要關係人：**\n{important_stakeholders}",
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

            # Third Slide: Key Stakeholders, Issue Gaps, and Future Service Opportunities
            
            issue_gaps = "\n".join(str(gap) for gap in trend["<g>缺口"])
            
            # 檢查 "<f>未來產品或服務機會點" 是否存在
            future_opportunities = "\n".join(str(opportunity) for opportunity in trend.get("<h>未來服務機會點", []))
            trend_summary = trend.get("<j>趨勢總結洞察", "")

            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}：{trend['<a>趨勢名稱']} ",
                content = f"**議題缺口：**\n{issue_gaps}\n\n**未來服務機會點：**\n{future_opportunities}\n\n**趨勢總結洞察：**\n{trend_summary}",
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

            # Fourth Slide: Key Drivers, Weak Signals, and Trend Summary Insights
            # key_drivers = "\n".join([f"{key}: {value}" for key, value in trend["<g>關鍵驅動因素"].items()])
            # key_drivers = "\n".join(["：".join([str(aspect), str(driver)])
            #                         for aspect, driver in trend["<g>關鍵驅動因素"].items()])
            # weak_signals = "\n".join([str(signal) for signal in trend["<h>微弱信號"]])
            # time_scale = trend.get("<i>時間尺度", "")
            
            # ExportManager.add_slide(prs = prs, topic = topic,
            #     title = f"{trend_type}{trend_number}：{trend['<a>趨勢名稱']} ",
            #     content = f"**趨勢總結洞察：**\n{trend_summary}",
            #     # content = f"**關鍵驅動因素：**\n{key_drivers}\n\n**微弱信號：**\n{weak_signals}\n\n**時間尺度：**{time_scale}\n\n**趨勢總結洞察：**\n{trend_summary}",
            #     title_font_size=32, content_font_size=28, rgb_color = rgb_color
            # )

    
    
    # Function that queries pics with project name as keyword, and downloads queried pictures
    @staticmethod
    def download_images(query, num_images):
        # * query
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": st.secrets["PEXEL_KEY"]}
        params = {"query": query, "per_page": num_images}
        response = requests.get(url, headers=headers, params=params)
        response.raise_for_status()
        st.write(response.content)
        results = response.json()
        image_urls = [photo["src"]["original"] for photo in results["photos"]]
        

        os.makedirs("images", exist_ok = True)
        for i, url in enumerate(image_urls):
            response = requests.get(url)
            with open(f"images/image_{i+1}.jpg", "wb") as f:
                f.write(response.content)
    
    # Function that add downloaded pics to ppt
    @staticmethod
    def add_image_slide(prs, image_file, title, title_font_size, background_color = RGBColor(240, 240, 240) , fixed_font_size=False, rgb_color = RGBColor(5, 5, 5),
                        ):
        slide_layout = prs.slide_layouts[0]  # Use title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title_placeholder = slide.shapes.title
        

        # Adjust title placeholder size to avoid overlapping with the black line
        title_placeholder.width = Inches(12.5)
        title_placeholder.height = Inches(1.2)
        title_placeholder.left = Inches(1)
        title_placeholder.top = Inches(0.4)

        # Set title and font size
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.clear()
        title_p = title_text_frame.add_paragraph()
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.bold = True
        title_p.font.name = 'Microsoft JhengHei'
        title_p.font.color.rgb = rgb_color
        title_p.alignment = PP_ALIGN.LEFT

        # Add a black line below the title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.9), Inches(14), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.width = Inches(0)

        # Add solid circle in the upper-right corner
        solid_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(13.7), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        solid_circle.fill.solid()
        solid_circle.fill.fore_color.rgb = rgb_color

        # Add hollow circle in the upper-right corner
        hollow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(14.5), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        hollow_circle.line.color.rgb = RGBColor(0, 0, 0)
        hollow_circle.line.width = Pt(1)
        hollow_circle.fill.background()

    

        # Add pics

        image_path = os.path.join("images", image_file)
        slide.shapes.add_picture(image_path, Inches(4), Inches(2), Inches(8), Inches(6))  # 調整圖片大小和位置

                
    class Export:
        
        @staticmethod
        def create_pptx(topic, data = None):

            '''convert JSON format report to PPT slides'''

            color = '#808080'
            rgb_color = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))

            prs = ExportManager.init_slides()
            assert data is not None, "Input data is required to create slides!"
            ExportManager.add_report_and_keywords(prs, data, topic, rgb_color)
            
            for i in range(1, 13):
                try:
                    ExportManager.add_trend_pages(prs, data, i, "趨勢", topic, rgb_color)
                except:
                    pass

            with open(f"output/{topic}.pptx", "wb") as f:
                prs.save(f)


            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            base64_pptx = base64.b64encode(buffer.read()).decode('utf-8')
            buffer.close()

            

            return base64_pptx

        # @staticmethod
        # def create_excel(start_date, end_date, topics: list):
            
        #     summary_reports = []

        #     buffer = io.BytesIO()
        #     with pd.ExcelWriter(buffer, engine = 'openpyxl') as writer:
        #         for topic in topics:
        #             try:
        #                 dfs_to_write = ExportManager.get_report_excels(start_date, end_date, topic)
        #                 summary_reports.append(dfs_to_write[1])
        #                 dfs_to_write[0].to_excel(writer, sheet_name = topic, index = False)
        #             except:
        #                 pass
        #         agg_res = pd.concat(summary_reports, ignore_index = False)
        #         agg_res.to_excel(writer, sheet_name = 'STEEP 趨勢報告與關鍵字', index = True)

        #     buffer.seek(0)
        #     b64_excel = base64.b64encode(buffer.read()).decode()
        #     buffer.close()
            

        #     return b64_excel
        
    